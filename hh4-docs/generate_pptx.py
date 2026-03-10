"""
generate_pptx.py
Generates HH4_Platform_Presentation.pptx — Splendora HH4 Digital Platform
Run: python generate_pptx.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Colour palette ──────────────────────────────────────────────────────────
C_NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
C_LIGHT  = RGBColor(0xF8, 0xF9, 0xFA)
C_GOLD   = RGBColor(0xC9, 0xA8, 0x4C)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LBLUE  = RGBColor(0xB0, 0xC4, 0xDE)
C_BLUE   = RGBColor(0x2E, 0x50, 0x90)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_GREEN  = RGBColor(0x00, 0xB0, 0x50)
C_RED    = RGBColor(0xFF, 0x40, 0x40)
C_ORANGE = RGBColor(0xFF, 0xA5, 0x00)

# ── Slide dimensions ────────────────────────────────────────────────────────
W = Cm(33.87)
H = Cm(19.05)


# ── Low-level helpers ────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color=C_WHITE, align=PP_ALIGN.LEFT,
                font_name="Calibri", wrap=True, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = word_wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_size=14, bold=False, italic=False,
                          color=C_WHITE, align=PP_ALIGN.LEFT,
                          font_name="Calibri", line_spacing=None):
    """lines: list of (text, size, bold, italic, color) tuples or plain strings."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            fs, fb, fi, fc = font_size, bold, italic, color
            text = item
        else:
            text = item[0]
            fs   = item[1] if len(item) > 1 else font_size
            fb   = item[2] if len(item) > 2 else bold
            fi   = item[3] if len(item) > 3 else italic
            fc   = item[4] if len(item) > 4 else color
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(fs)
        run.font.bold = fb
        run.font.italic = fi
        run.font.color.rgb = fc
    return txb


def add_gold_line(slide, y=None):
    """Thin gold horizontal separator."""
    if y is None:
        y = H - Cm(1.4)
    from pptx.util import Pt as UPt
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Cm(1), y, W - Cm(1), y
    )
    line.line.color.rgb = C_GOLD
    line.line.width = Pt(1.5)


def add_logo(slide, dark=True):
    """HH4 text logo top-right."""
    color = C_GOLD
    add_textbox(slide, "HH4",
                W - Cm(3.2), Cm(0.3), Cm(2.8), Cm(0.9),
                font_size=16, bold=True, color=color,
                align=PP_ALIGN.RIGHT)


def add_slide_number(slide, num, dark=True):
    color = C_GOLD if dark else C_BLUE
    add_textbox(slide, str(num),
                W - Cm(2), H - Cm(1.1), Cm(1.6), Cm(0.8),
                font_size=10, bold=False, color=color,
                align=PP_ALIGN.RIGHT)


def chrome(slide, num, dark=True):
    """Add logo + gold line + slide number."""
    add_logo(slide, dark)
    add_gold_line(slide)
    add_slide_number(slide, num, dark)


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=Pt(1)):
    from pptx.util import Pt as UPt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def box_with_text(slide, left, top, width, height,
                  text, font_size=14, bold=False, italic=False,
                  text_color=C_WHITE, fill_color=C_BLUE,
                  border_color=C_GOLD, border_width=Pt(1.5),
                  align=PP_ALIGN.CENTER, font_name="Calibri"):
    add_rect(slide, left, top, width, height,
             fill_color=fill_color,
             line_color=border_color, line_width=border_width)
    add_textbox(slide, text, left + Cm(0.2), top + Cm(0.15),
                width - Cm(0.4), height - Cm(0.3),
                font_size=font_size, bold=bold, italic=italic,
                color=text_color, align=align, font_name=font_name)


def section_title(slide, text, dark=True):
    color = C_GOLD if dark else C_BLUE
    add_textbox(slide, text,
                Cm(1.2), Cm(0.9), W - Cm(5), Cm(1.2),
                font_size=24, bold=True, color=color,
                align=PP_ALIGN.LEFT)


def table_cell_set(cell, text, font_size=11, bold=False,
                   text_color=C_BLACK, fill_color=None,
                   align=PP_ALIGN.CENTER, font_name="Calibri"):
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    if fill_color:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(
            fill_color[0], fill_color[1], fill_color[2]))


# ── Slide builders ───────────────────────────────────────────────────────────

def slide01(prs):
    """COVER"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)

    # Gold line divider mid-section
    div_y = Cm(7.5)
    line = sl.shapes.add_connector(1, Cm(2), div_y, W - Cm(2), div_y)
    line.line.color.rgb = C_GOLD
    line.line.width = Pt(2)

    add_textbox(sl, "SPLENDORA HH4 DIGITAL PLATFORM",
                Cm(2), Cm(2.5), W - Cm(4), Cm(2.8),
                font_size=40, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)

    add_textbox(sl, "Where Luxury Living Begins Before Move-In",
                Cm(2), Cm(5.5), W - Cm(4), Cm(1.4),
                font_size=22, italic=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(sl, "Investment Proposal  |  >$10,000,000 USD",
                Cm(2), Cm(8.3), W - Cm(4), Cm(1.1),
                font_size=18, bold=False, color=C_LBLUE,
                align=PP_ALIGN.CENTER)

    add_textbox(sl, "March 2026  |  Strictly Confidential",
                Cm(2), H - Cm(2.2), W - Cm(4), Cm(0.9),
                font_size=13, italic=True, color=C_LBLUE,
                align=PP_ALIGN.CENTER)

    chrome(sl, 1, dark=True)


def slide02(prs):
    """EXECUTIVE HOOK"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "3 Numbers That Define This Investment")

    stats = [
        ("$250,000,000", "Total GTV to Protect & Unlock"),
        ("$11,900",      "HNW Prospect's Time Value Per Hour"),
        ("87%",          "Buyers Research Digitally Before Physical Visit"),
    ]
    box_w = Cm(9.2)
    box_h = Cm(6.5)
    top   = Cm(3.2)
    for i, (val, label) in enumerate(stats):
        lft = Cm(1.2) + i * (box_w + Cm(0.8))
        add_rect(sl, lft, top, box_w, box_h,
                 fill_color=None, line_color=C_GOLD, line_width=Pt(2))
        add_textbox(sl, val,
                    lft + Cm(0.3), top + Cm(0.7), box_w - Cm(0.6), Cm(2.4),
                    font_size=34, bold=True, color=C_GOLD,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, label,
                    lft + Cm(0.3), top + Cm(3.3), box_w - Cm(0.6), Cm(2.8),
                    font_size=14, italic=False, color=C_LBLUE,
                    align=PP_ALIGN.CENTER)

    add_textbox(sl,
                "Every second of experience counts. Every data point matters.",
                Cm(2), H - Cm(2.4), W - Cm(4), Cm(0.9),
                font_size=12, italic=True, color=C_LBLUE,
                align=PP_ALIGN.CENTER)

    chrome(sl, 2, dark=True)


def slide03(prs):
    """THE PROBLEM"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "Today's Real Estate Sales Process Is Failing Premium Buyers")

    pains = [
        "❌  Brochure PDF — No tracking, no personalization",
        "❌  Physical showroom — Geography-limited, 9–6 only",
        "❌  Sales call — Inconsistent, human-dependent",
        "❌  Post-sale — Zero ecosystem. Zero loyalty.",
    ]
    box_h = Cm(2.0)
    for i, pain in enumerate(pains):
        top = Cm(3.0) + i * (box_h + Cm(0.35))
        add_rect(sl, Cm(1.5), top, W - Cm(3), box_h,
                 fill_color=RGBColor(0x14, 0x2A, 0x42),
                 line_color=C_GOLD, line_width=Pt(1))
        add_textbox(sl, pain,
                    Cm(2.0), top + Cm(0.25), W - Cm(4), box_h - Cm(0.3),
                    font_size=16, color=C_WHITE)

    # Gold callout
    callout_top = Cm(13.7)
    add_rect(sl, Cm(1.5), callout_top, W - Cm(3), Cm(1.8),
             fill_color=C_GOLD, line_color=None)
    add_textbox(sl,
                "HNW/UHNW buyers expect Amazon-level personalization "
                "with Ritz-Carlton-level service.",
                Cm(2.0), callout_top + Cm(0.2), W - Cm(4), Cm(1.4),
                font_size=15, bold=True, italic=True,
                color=C_NAVY, align=PP_ALIGN.CENTER)

    chrome(sl, 3, dark=True)


def slide04(prs):
    """THE MARKET MOMENT"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "Why Now? Three Converging Forces")

    cols = [
        ("🌐 Digital-native HNW",
         "78% of Vietnam UHNW under 50 prefer digital-first engagement"),
        ("📊 First-mover Window",
         "No luxury developer in Vietnam has deployed a full-lifecycle "
         "digital platform at this scale"),
        ("🏆 Data Is the New Land Bank",
         "Whoever owns resident data owns 20+ years of recurring revenue"),
    ]
    col_w = Cm(9.5)
    col_h = Cm(9.0)
    top   = Cm(3.2)
    for i, (title, body) in enumerate(cols):
        lft = Cm(1.2) + i * (col_w + Cm(0.6))
        add_rect(sl, lft, top, col_w, col_h,
                 fill_color=C_BLUE,
                 line_color=C_GOLD, line_width=Pt(1.5))
        add_textbox(sl, title,
                    lft + Cm(0.3), top + Cm(0.5), col_w - Cm(0.6), Cm(2.0),
                    font_size=16, bold=True, color=C_GOLD,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, body,
                    lft + Cm(0.4), top + Cm(2.8), col_w - Cm(0.8), Cm(5.5),
                    font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

    chrome(sl, 4, dark=True)


def slide05(prs):
    """OUR ANSWER"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "Introducing the Splendora HH4 Digital Platform")

    box_w = Cm(12.5)
    box_h = Cm(5.5)
    top   = Cm(3.5)

    add_rect(sl, Cm(1.2), top, box_w, box_h,
             fill_color=C_BLUE, line_color=C_GOLD, line_width=Pt(2))
    add_textbox(sl, "🏛️ IVA METAVERSE SHOWROOM",
                Cm(1.4), top + Cm(0.4), box_w - Cm(0.4), Cm(1.4),
                font_size=17, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)
    add_textbox(sl, "Pre-sale & Sale",
                Cm(1.4), top + Cm(1.9), box_w - Cm(0.4), Cm(1.0),
                font_size=15, italic=True, color=C_LBLUE,
                align=PP_ALIGN.CENTER)
    add_textbox(sl, "3D tours · AI scoring · Live sessions\nUnit configurator · Auto proposals",
                Cm(1.4), top + Cm(3.0), box_w - Cm(0.4), Cm(2.0),
                font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Arrow
    add_textbox(sl, "◄────────────►",
                Cm(13.9), top + Cm(2.2), Cm(2.8), Cm(1.0),
                font_size=18, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)

    add_rect(sl, Cm(16.8), top, box_w, box_h,
             fill_color=C_BLUE, line_color=C_GOLD, line_width=Pt(2))
    add_textbox(sl, "🏠 CONCIERGE 360°",
                Cm(17.0), top + Cm(0.4), box_w - Cm(0.4), Cm(1.4),
                font_size=17, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)
    add_textbox(sl, "Post-sale & Lifetime",
                Cm(17.0), top + Cm(1.9), box_w - Cm(0.4), Cm(1.0),
                font_size=15, italic=True, color=C_LBLUE,
                align=PP_ALIGN.CENTER)
    add_textbox(sl, "AI concierge · Smart home · Services\nCommunity · Loyalty & rewards",
                Cm(17.0), top + Cm(3.0), box_w - Cm(0.4), Cm(2.0),
                font_size=13, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(sl, "One Platform. Full Lifecycle. Unmatched Experience.",
                Cm(2), Cm(10.2), W - Cm(4), Cm(1.0),
                font_size=18, bold=True, italic=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)

    tags = ["[3D/VR]", "[AI Scoring]", "[Live Tours]",
            "[Smart Home]", "[Partner Ecosystem]"]
    tag_w = Cm(5.4)
    tag_top = Cm(11.6)
    for i, tag in enumerate(tags):
        lft = Cm(1.2) + i * (tag_w + Cm(0.5))
        add_rect(sl, lft, tag_top, tag_w, Cm(0.85),
                 fill_color=C_GOLD, line_color=None)
        add_textbox(sl, tag, lft + Cm(0.1), tag_top + Cm(0.05),
                    tag_w - Cm(0.2), Cm(0.75),
                    font_size=13, bold=True, color=C_NAVY,
                    align=PP_ALIGN.CENTER)

    chrome(sl, 5, dark=True)


def slide06(prs):
    """PLATFORM OVERVIEW"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "One Platform. Two Worlds. Infinite Touchpoints.")

    modules = [
        ("1. Marketing Module",      "Attract & Generate Quality Leads"),
        ("2. Experience Module",     "Immersive Digital Experiences"),
        ("3. Sales Enablement",      "Shorten Sales Cycle, Close Faster"),
        ("4. Concierge & Service",   "Personalize & Retain for Life"),
        ("5. Admin & Operations",    "Control, Monitor, Optimize"),
    ]
    row_h = Cm(1.9)
    for i, (mod, desc) in enumerate(modules):
        top = Cm(2.8) + i * (row_h + Cm(0.2))
        add_rect(sl, Cm(1.2), top, W - Cm(2.4), row_h,
                 fill_color=C_BLUE,
                 line_color=C_GOLD, line_width=Pt(1))
        add_textbox(sl, mod,
                    Cm(1.6), top + Cm(0.2), Cm(9), row_h - Cm(0.2),
                    font_size=15, bold=True, color=C_WHITE)
        add_textbox(sl, desc,
                    Cm(11), top + Cm(0.2), Cm(20), row_h - Cm(0.2),
                    font_size=14, color=C_LBLUE)

    add_textbox(sl,
                "User Roles:  Guest  ·  Registered  ·  VIP  ·  "
                "Sales / Marketing  ·  Admin",
                Cm(1.2), H - Cm(2.6), W - Cm(2.4), Cm(0.9),
                font_size=13, bold=True, color=C_BLUE,
                align=PP_ALIGN.CENTER)

    chrome(sl, 6, dark=False)


def slide07(prs):
    """IVA METAVERSE SHOWROOM"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "Your Showroom Is Now Open. 24/7. Worldwide.")

    cards = [
        ("🌐 Metaverse Showroom",
         "Walk through HH4 before it's built.\n"
         "3D masterplan, unit configurator,\nlifestyle zones."),
        ("🤖 AI Behavioral Scoring",
         "Know who's serious before they call.\n"
         "Auto-route hot leads in <5 minutes."),
        ("🎯 Guided Tour Engine",
         "Sales-led live tours.\nAvatar + Voice/Video.\n"
         "Co-navigation in real-time."),
    ]
    card_w = Cm(9.8)
    card_h = Cm(8.2)
    top    = Cm(3.0)
    for i, (title, body) in enumerate(cards):
        lft = Cm(1.2) + i * (card_w + Cm(0.55))
        add_rect(sl, lft, top, card_w, card_h,
                 fill_color=RGBColor(0x14, 0x2A, 0x42),
                 line_color=C_GOLD, line_width=Pt(2))
        add_textbox(sl, title,
                    lft + Cm(0.4), top + Cm(0.5), card_w - Cm(0.8), Cm(1.5),
                    font_size=16, bold=True, color=C_GOLD,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, body,
                    lft + Cm(0.4), top + Cm(2.3), card_w - Cm(0.8), Cm(5.4),
                    font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(sl, Cm(12.5), H - Cm(2.6), Cm(8.8), Cm(0.9),
             fill_color=C_GOLD, line_color=None)
    add_textbox(sl, "[LIVE DEMO — 60 seconds]",
                Cm(12.7), H - Cm(2.55), Cm(8.4), Cm(0.8),
                font_size=14, bold=True, color=C_NAVY,
                align=PP_ALIGN.CENTER)

    chrome(sl, 7, dark=True)


def slide08(prs):
    """CUSTOMER JOURNEY: IVA"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "From First Click to Signed Contract in 45 Days")

    stages = [
        ("1\nAwareness",    "Digital ads\n& social",  "60%"),
        ("2\nSelf-guided",  "3D IVA\nexploration",    "45%"),
        ("3\nQualification","AI scoring\n& routing",  "35%"),
        ("4\nPrivate\nViewing","Live guided\ntour",   "70%"),
        ("5\nO2O",          "Site visit\n& meeting",  "80%"),
        ("6\nDecision",     "Proposal\n& deposit",    "25%"),
    ]
    box_w = Cm(4.9)
    box_h = Cm(6.0)
    top   = Cm(3.2)
    for i, (stage, feat, cvr) in enumerate(stages):
        lft = Cm(0.9) + i * (box_w + Cm(0.38))
        color = C_GOLD if i % 2 == 0 else C_BLUE
        add_rect(sl, lft, top, box_w, box_h,
                 fill_color=color,
                 line_color=None)
        add_textbox(sl, stage,
                    lft + Cm(0.2), top + Cm(0.4), box_w - Cm(0.4), Cm(2.0),
                    font_size=13, bold=True,
                    color=C_NAVY if i % 2 == 0 else C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, feat,
                    lft + Cm(0.2), top + Cm(2.5), box_w - Cm(0.4), Cm(2.0),
                    font_size=12,
                    color=C_NAVY if i % 2 == 0 else C_LBLUE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, f"Conv.\n{cvr}",
                    lft + Cm(0.2), top + Cm(4.5), box_w - Cm(0.4), Cm(1.2),
                    font_size=11, bold=True,
                    color=C_NAVY if i % 2 == 0 else C_GOLD,
                    align=PP_ALIGN.CENTER)

    add_textbox(sl,
                "Each stage captures behavioural data that feeds AI scoring — "
                "shortening the full cycle to ≤45 days.",
                Cm(1.2), H - Cm(2.5), W - Cm(2.4), Cm(0.9),
                font_size=12, italic=True, color=C_BLUE,
                align=PP_ALIGN.CENTER)
    chrome(sl, 8, dark=False)


def slide09(prs):
    """AI BEHAVIORAL SCORING"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "We Know Intent Before Your Customer Does")

    # Left — INPUT SIGNALS
    signals = [
        "Time in 3D unit view",
        "Material / config choices",
        "Shortlist additions",
        "Return visit frequency",
        "Video engagement depth",
        "Form completion rate",
    ]
    add_textbox(sl, "INPUT SIGNALS",
                Cm(0.8), Cm(3.0), Cm(8), Cm(1.1),
                font_size=14, bold=True, color=C_GOLD)
    for i, sig in enumerate(signals):
        add_textbox(sl, f"▸  {sig}",
                    Cm(1.0), Cm(4.2) + i * Cm(1.3), Cm(7.5), Cm(1.1),
                    font_size=13, color=C_LBLUE)

    # Centre — AI ENGINE circle (rectangle approximation)
    add_rect(sl, Cm(10.3), Cm(5.2), Cm(4.2), Cm(4.2),
             fill_color=C_GOLD, line_color=C_WHITE, line_width=Pt(2))
    add_textbox(sl, "AI\nENGINE",
                Cm(10.5), Cm(6.2), Cm(3.8), Cm(2.2),
                font_size=18, bold=True, color=C_NAVY,
                align=PP_ALIGN.CENTER)

    # Right — OUTPUT
    outputs = [
        ("Score >75", "HOT LEAD", "→ Senior Sales (5 min alert)", C_RED),
        ("Score 45–74", "WARM", "→ Concierge Nurture", C_ORANGE),
        ("Score <45", "COLD", "→ Automated Sequence", C_LBLUE),
    ]
    add_textbox(sl, "OUTPUT",
                Cm(16.0), Cm(3.0), Cm(12), Cm(1.1),
                font_size=14, bold=True, color=C_GOLD)
    for i, (score, label, action, col) in enumerate(outputs):
        top = Cm(4.2) + i * Cm(2.8)
        add_rect(sl, Cm(16.0), top, Cm(13.5), Cm(2.4),
                 fill_color=RGBColor(0x14, 0x2A, 0x42),
                 line_color=col, line_width=Pt(1.5))
        add_textbox(sl, f"{score}  →  {label}",
                    Cm(16.3), top + Cm(0.2), Cm(13), Cm(1.0),
                    font_size=14, bold=True, color=col)
        add_textbox(sl, action,
                    Cm(16.3), top + Cm(1.1), Cm(13), Cm(1.0),
                    font_size=12, italic=True, color=C_LBLUE)

    add_textbox(sl,
                "Example: 8 min in Hero Unit + Material Config + Shortlist  "
                "=  Score 89/100  →  HOT",
                Cm(1.2), H - Cm(2.5), W - Cm(2.4), Cm(0.9),
                font_size=12, italic=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)
    chrome(sl, 9, dark=True)


def slide10(prs):
    """CONCIERGE 360°"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "The World's Most Personalized Home Assistant")

    pillars = [
        ("🤖 AI Personal Concierge",
         "Not a chatbot. A lifestyle curator.\nLearns, predicts, acts."),
        ("🏪 Service Ecosystem",
         "500+ services, 1 app, zero friction.\nBuilding + Partners + Smart Home."),
        ("👥 Community & Loyalty",
         "Residents who stay, spend, and refer.\nGamification + Rewards."),
    ]
    card_w = Cm(9.8)
    card_h = Cm(7.5)
    top    = Cm(3.2)
    for i, (title, body) in enumerate(pillars):
        lft = Cm(1.2) + i * (card_w + Cm(0.55))
        add_rect(sl, lft, top, card_w, card_h,
                 fill_color=C_BLUE,
                 line_color=C_GOLD, line_width=Pt(2))
        add_textbox(sl, title,
                    lft + Cm(0.4), top + Cm(0.6), card_w - Cm(0.8), Cm(1.6),
                    font_size=16, bold=True, color=C_GOLD,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, body,
                    lft + Cm(0.4), top + Cm(2.5), card_w - Cm(0.8), Cm(4.5),
                    font_size=14, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(sl, Cm(12.5), H - Cm(2.65), Cm(8.8), Cm(0.9),
             fill_color=C_GOLD, line_color=None)
    add_textbox(sl, "[APP DEMO — 45 seconds]",
                Cm(12.7), H - Cm(2.6), Cm(8.4), Cm(0.8),
                font_size=14, bold=True, color=C_NAVY,
                align=PP_ALIGN.CENTER)

    chrome(sl, 10, dark=True)


def slide11(prs):
    """RESIDENT JOURNEY"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "From Move-In Day to Lifetime Advocate")

    stages = [
        ("1\nAwareness",   "Marketing\ntouch",       "Curious"),
        ("2\nOnboarding",  "App setup\nSmart home",  "Excited"),
        ("3\nDaily\nLiving","Services\n& requests",  "Satisfied"),
        ("4\nCommunity",   "Events\n& social",       "Connected"),
        ("5\nExpansion",   "Upgrade\n& referral",    "Loyal"),
        ("6\nLifetime\nAdvocate","Repeat &\nrefer",  "Advocate"),
    ]
    box_w = Cm(4.9)
    box_h = Cm(6.5)
    top   = Cm(3.2)
    for i, (stage, kpi, emotion) in enumerate(stages):
        lft = Cm(0.9) + i * (box_w + Cm(0.38))
        color = C_GOLD if i % 2 == 0 else C_BLUE
        add_rect(sl, lft, top, box_w, box_h,
                 fill_color=color, line_color=None)
        add_textbox(sl, stage,
                    lft + Cm(0.2), top + Cm(0.35), box_w - Cm(0.4), Cm(2.0),
                    font_size=13, bold=True,
                    color=C_NAVY if i % 2 == 0 else C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, kpi,
                    lft + Cm(0.2), top + Cm(2.5), box_w - Cm(0.4), Cm(2.0),
                    font_size=12,
                    color=C_NAVY if i % 2 == 0 else C_LBLUE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, emotion,
                    lft + Cm(0.2), top + Cm(4.8), box_w - Cm(0.4), Cm(1.2),
                    font_size=11, italic=True, bold=True,
                    color=C_NAVY if i % 2 == 0 else C_GOLD,
                    align=PP_ALIGN.CENTER)

    add_textbox(sl, "CLV grows 40%+ when residents engage with Concierge ≥3×/week",
                Cm(1.2), H - Cm(2.5), W - Cm(2.4), Cm(0.9),
                font_size=12, italic=True, color=C_BLUE,
                align=PP_ALIGN.CENTER)
    chrome(sl, 11, dark=False)


def slide12(prs):
    """AI CONCIERGE IN ACTION"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "Meet Your Residents' 24/7 Personal Concierge")

    convos = [
        ("A", '"Book a table at rooftop for 4 tonight"',
               '"Done in 8 seconds ✓"'),
        ("B", '"I need a PT tomorrow morning"',
               '"Booked 7AM, added to calendar ✓"'),
        ("C", '"Dinner party Saturday"',
               '"Chef + florist + wine, one tap ✓"'),
    ]
    box_h = Cm(4.0)
    for i, (lbl, req, resp) in enumerate(convos):
        top = Cm(3.2) + i * (box_h + Cm(0.4))
        add_rect(sl, Cm(1.2), top, W - Cm(2.4), box_h,
                 fill_color=RGBColor(0x14, 0x2A, 0x42),
                 line_color=C_GOLD, line_width=Pt(1.5))
        add_textbox(sl, lbl,
                    Cm(1.5), top + Cm(0.3), Cm(1.2), Cm(3.4),
                    font_size=18, bold=True, color=C_GOLD,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, f"Resident:  {req}",
                    Cm(3.2), top + Cm(0.4), W - Cm(6), Cm(1.4),
                    font_size=14, italic=True, color=C_LBLUE)
        add_textbox(sl, f"AI:  {resp}",
                    Cm(3.2), top + Cm(1.9), W - Cm(6), Cm(1.4),
                    font_size=14, bold=True, color=C_GREEN)

    add_textbox(sl,
                "Response time: <1.5 seconds  |  "
                "Satisfaction score: >4.2 / 5.0",
                Cm(2), H - Cm(2.5), W - Cm(4), Cm(0.9),
                font_size=13, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)
    chrome(sl, 12, dark=True)


def slide13(prs):
    """DATA INTELLIGENCE"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "The Platform That Gets Smarter Every Day")

    cols_data = [
        ("COLLECT", C_BLUE, [
            "Behavioral (IVA + App)",
            "Transaction data",
            "Smart Home telemetry",
            "Engagement events",
        ]),
        ("PROCESS", C_GOLD, [
            "CDP Unified Profile",
            "ML Pipeline",
            "Real-time Events",
            "Data Warehouse",
        ]),
        ("ACTIVATE", C_BLUE, [
            "Personalised Content",
            "Lead Routing",
            "Triggered Campaigns",
            "BI Dashboards",
        ]),
    ]
    col_w = Cm(8.8)
    col_h = Cm(9.5)
    top   = Cm(3.0)
    for i, (hdr, col, items) in enumerate(cols_data):
        lft = Cm(1.8) + i * (col_w + Cm(1.4))
        if i < 2:
            arr_lft = lft + col_w + Cm(0.2)
            add_textbox(sl, "▶▶",
                        arr_lft, top + Cm(4.0), Cm(1.0), Cm(1.2),
                        font_size=20, bold=True, color=C_GOLD,
                        align=PP_ALIGN.CENTER)
        add_rect(sl, lft, top, col_w, col_h,
                 fill_color=col if col == C_GOLD else RGBColor(0xE8, 0xF0, 0xFE),
                 line_color=col, line_width=Pt(2))
        add_textbox(sl, hdr,
                    lft + Cm(0.3), top + Cm(0.5), col_w - Cm(0.6), Cm(1.2),
                    font_size=17, bold=True,
                    color=C_NAVY if col == C_GOLD else C_BLUE,
                    align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            add_textbox(sl, f"▸  {item}",
                        lft + Cm(0.5), top + Cm(2.0) + j * Cm(1.65),
                        col_w - Cm(1.0), Cm(1.5),
                        font_size=13,
                        color=C_NAVY if col == C_GOLD else C_BLACK)

    add_textbox(sl,
                "Single source of truth. Predictive, not reactive.",
                Cm(1.2), H - Cm(2.5), W - Cm(2.4), Cm(0.9),
                font_size=13, bold=True, italic=True, color=C_BLUE,
                align=PP_ALIGN.CENTER)
    chrome(sl, 13, dark=False)


def slide14(prs):
    """TECHNICAL CREDIBILITY"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "Built for Scale. Designed for Security. Ready for Growth.")

    quads = [
        ("☁️ Cloud-Native",
         "AWS/GCP · Kubernetes\n99.9% SLA · Microservices\nAuto-scaling · CI/CD",
         Cm(1.2), Cm(3.2)),
        ("🔒 Security & Compliance",
         "PDPA · ISO 27001 target\nWAF · 2FA · E2E Encryption\nAnnual Pen Testing",
         W / 2 + Cm(0.3), Cm(3.2)),
        ("🔗 Integration-Ready",
         "CRM / ERP / Smart Home\nPayment · Marketing Auto\nAPI-first · Webhooks",
         Cm(1.2), Cm(10.0)),
        ("📈 Scalable",
         "Replicate to next project\nat 30% incremental cost\nAPI-first architecture",
         W / 2 + Cm(0.3), Cm(10.0)),
    ]
    quad_w = W / 2 - Cm(1.5)
    quad_h = Cm(5.8)
    for title, body, lft, top in quads:
        add_rect(sl, lft, top, quad_w, quad_h,
                 fill_color=C_BLUE,
                 line_color=C_GOLD, line_width=Pt(1.5))
        add_textbox(sl, title,
                    lft + Cm(0.4), top + Cm(0.4), quad_w - Cm(0.8), Cm(1.4),
                    font_size=16, bold=True, color=C_GOLD)
        add_textbox(sl, body,
                    lft + Cm(0.4), top + Cm(2.0), quad_w - Cm(0.8), Cm(3.4),
                    font_size=13, color=C_WHITE)

    chrome(sl, 14, dark=False)


def slide15(prs):
    """THE FINANCIAL CASE"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "$11M Investment. $21.9M Net Return in 5 Years.")

    years   = ["Y1", "Y2", "Y3", "Y4", "Y5"]
    values  = [-1.0, 9.5, 10.4, 1.3, 1.7]
    bar_w   = Cm(4.5)
    max_val = 12.0
    bar_area_h = Cm(8.0)
    zero_y  = Cm(12.8)
    base_x  = Cm(2.0)

    for i, (yr, val) in enumerate(zip(years, values)):
        lft = base_x + i * (bar_w + Cm(0.8))
        bar_h = abs(val) / max_val * bar_area_h
        if val >= 0:
            top = zero_y - bar_h
            col = C_BLUE
        else:
            top = zero_y
            col = C_RED

        add_rect(sl, lft, top, bar_w, bar_h,
                 fill_color=col, line_color=None)
        label = f"${val:+.1f}M"
        label_top = top - Cm(0.8) if val >= 0 else top + bar_h + Cm(0.1)
        add_textbox(sl, label,
                    lft, label_top, bar_w, Cm(0.7),
                    font_size=13, bold=True,
                    color=C_BLUE if val >= 0 else C_RED,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, yr,
                    lft, zero_y + Cm(0.2), bar_w, Cm(0.7),
                    font_size=13, bold=True, color=C_BLACK,
                    align=PP_ALIGN.CENTER)

    # Zero line
    line = sl.shapes.add_connector(1, base_x, zero_y, base_x + Cm(28), zero_y)
    line.line.color.rgb = C_BLACK
    line.line.width = Pt(1)

    add_rect(sl, Cm(7.5), H - Cm(2.7), Cm(18), Cm(1.1),
             fill_color=C_GOLD, line_color=None)
    add_textbox(sl,
                "Break-even: Month 18  |  5-Year ROI: 199%  |  "
                "Bear scenario: break-even by Year 2",
                Cm(7.7), H - Cm(2.65), Cm(17.6), Cm(0.95),
                font_size=13, bold=True, color=C_NAVY,
                align=PP_ALIGN.CENTER)

    chrome(sl, 15, dark=False)


def slide16(prs):
    """ROI BREAKDOWN"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "Three Ways This Platform Pays for Itself")

    cards = [
        ("💰 SALES ACCELERATION",
         "+3.5% conversion delta",
         "$12.6M additional\ndeveloper margin/year"),
        ("🔄 CONCIERGE REVENUE",
         "Recurring from Year 2",
         "$985K–$1.1M/year,\ngrowing annually"),
        ("✨ BRAND PREMIUM",
         "+15–20% unit pricing",
         "$5.6M–$7.5M uplift\non GTV"),
    ]
    card_w = Cm(9.8)
    card_h = Cm(8.0)
    top    = Cm(3.2)
    for i, (title, sub, val) in enumerate(cards):
        lft = Cm(1.2) + i * (card_w + Cm(0.55))
        add_rect(sl, lft, top, card_w, card_h,
                 fill_color=C_BLUE,
                 line_color=C_GOLD, line_width=Pt(2))
        add_textbox(sl, title,
                    lft + Cm(0.4), top + Cm(0.5), card_w - Cm(0.8), Cm(1.6),
                    font_size=15, bold=True, color=C_GOLD,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, sub,
                    lft + Cm(0.4), top + Cm(2.5), card_w - Cm(0.8), Cm(1.4),
                    font_size=14, italic=True, color=C_LBLUE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, val,
                    lft + Cm(0.4), top + Cm(4.3), card_w - Cm(0.8), Cm(3.0),
                    font_size=16, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)

    add_rect(sl, Cm(5.5), H - Cm(2.65), Cm(22), Cm(1.1),
             fill_color=C_GOLD, line_color=None)
    add_textbox(sl, "Total 5-Year Net Return: $21.9M on $11M investment",
                Cm(5.7), H - Cm(2.6), Cm(21.6), Cm(0.95),
                font_size=15, bold=True, color=C_NAVY,
                align=PP_ALIGN.CENTER)

    chrome(sl, 16, dark=True)


def slide17(prs):
    """COST BREAKDOWN"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "How the $11M Is Deployed")

    segments = [
        ("3D / Metaverse",  "28%",  Cm(1.2)),
        ("Backend Infra",   "22%",  Cm(5.5)),
        ("AI / ML",         "15%",  Cm(9.8)),
        ("Mobile App",      "10%",  Cm(14.1)),
        ("Integrations",    "10%",  Cm(17.2)),
        ("Cloud & DevOps",  "8%",   Cm(20.3)),
        ("Analytics",       "5%",   Cm(23.4)),
        ("QA & Security",   "2%",   Cm(26.5)),
    ]
    bar_top = Cm(3.2)
    bar_h   = Cm(3.2)
    colours = [C_GOLD, C_BLUE, RGBColor(0x4A, 0x90, 0xD9),
               RGBColor(0x7B, 0x68, 0xEE), RGBColor(0x00, 0xC8, 0xA0),
               RGBColor(0xFF, 0xA5, 0x00), RGBColor(0xFF, 0x6B, 0x6B),
               RGBColor(0x98, 0xFB, 0x98)]
    seg_w_map = [6.0, 4.7, 3.2, 2.2, 2.2, 1.7, 1.1, 0.5]
    for i, ((label, pct, _lft), col, sw) in \
            enumerate(zip(segments, colours, seg_w_map)):
        lft = Cm(0.8) + sum(Cm(seg_w_map[j] + 0.15) for j in range(i))
        bw = Cm(sw)
        add_rect(sl, lft, bar_top, bw, bar_h,
                 fill_color=col, line_color=C_WHITE, line_width=Pt(0.5))
        if sw > 1.0:
            add_textbox(sl, pct,
                        lft + Cm(0.1), bar_top + Cm(0.2), bw - Cm(0.2), Cm(0.9),
                        font_size=12, bold=True,
                        color=C_NAVY if col == C_GOLD else C_WHITE,
                        align=PP_ALIGN.CENTER)
            add_textbox(sl, label,
                        lft + Cm(0.1), bar_top + Cm(1.2), bw - Cm(0.2), Cm(1.8),
                        font_size=10,
                        color=C_NAVY if col == C_GOLD else C_WHITE,
                        align=PP_ALIGN.CENTER)

    # Phase timeline
    phases = [
        ("Phase 1  $3.5M", "M1 – M8",   "Foundation: IVA MVP · Admin · CRM"),
        ("Phase 2  $3.5M", "M9 – M16",  "Intelligence: AI Scoring · Concierge · Payment"),
        ("Phase 3  $3.0M", "M17 – M24", "Ecosystem: Full AI · Partners · Community"),
    ]
    ph_top = Cm(8.0)
    ph_w   = (W - Cm(2.4)) / 3
    for i, (ph, period, desc) in enumerate(phases):
        lft = Cm(1.2) + i * ph_w
        add_rect(sl, lft, ph_top, ph_w - Cm(0.2), Cm(5.5),
                 fill_color=C_BLUE, line_color=C_GOLD, line_width=Pt(1.5))
        add_textbox(sl, ph,
                    lft + Cm(0.3), ph_top + Cm(0.4),
                    ph_w - Cm(0.8), Cm(1.2),
                    font_size=14, bold=True, color=C_GOLD)
        add_textbox(sl, period,
                    lft + Cm(0.3), ph_top + Cm(1.7),
                    ph_w - Cm(0.8), Cm(0.9),
                    font_size=13, bold=True, color=C_WHITE)
        add_textbox(sl, desc,
                    lft + Cm(0.3), ph_top + Cm(2.8),
                    ph_w - Cm(0.8), Cm(2.3),
                    font_size=12, color=C_LBLUE)

    add_textbox(sl,
                "Phased investment. Milestone-gated. No large upfront risk.",
                Cm(1.2), H - Cm(2.5), W - Cm(2.4), Cm(0.9),
                font_size=12, bold=True, italic=True, color=C_BLUE,
                align=PP_ALIGN.CENTER)
    chrome(sl, 17, dark=False)


def slide18(prs):
    """COMPETITIVE ADVANTAGE"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl,
                  "This Has Never Been Done in Vietnam Luxury Real Estate")

    rows_data = [
        ("Feature",                "HH4 Platform", "Competitor A", "Market Average"),
        ("Metaverse Showroom",     "✅ Full",       "⚡ Limited",   "❌"),
        ("AI Behavioral Scoring",  "✅",            "❌",           "❌"),
        ("Post-sale Ecosystem",    "✅ Full",       "❌",           "❌"),
        ("Unified Data Platform",  "✅",            "❌",           "⚡ Limited"),
        ("Partner Marketplace",    "✅",            "❌",           "❌"),
        ("O2O Journey Bridge",     "✅",            "❌",           "❌"),
        ("Multi-project Scalable", "✅",            "❌",           "⚡"),
        ("Smart Home Integration", "✅",            "❌",           "❌"),
        ("AI Concierge",           "✅",            "❌",           "❌"),
        ("Live Guided Tour",       "✅",            "⚡",           "❌"),
        ("Auto-Proposal Gen",      "✅",            "❌",           "❌"),
        ("Loyalty Ecosystem",      "✅",            "❌",           "❌"),
    ]
    num_rows = len(rows_data)
    num_cols = 4
    tbl_top  = Cm(2.6)
    tbl_h    = H - tbl_top - Cm(2.5)
    tbl_w    = W - Cm(2.4)
    tbl = sl.shapes.add_table(num_rows, num_cols,
                               Cm(1.2), tbl_top, tbl_w, tbl_h).table
    col_widths = [Cm(10), Cm(7), Cm(7), Cm(7)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            if ri == 0:
                table_cell_set(cell, cell_text, font_size=12, bold=True,
                               text_color=C_WHITE,
                               fill_color=C_BLUE,
                               align=PP_ALIGN.CENTER)
            else:
                if ci == 0:
                    fill = RGBColor(0xF0, 0xF4, 0xFF)
                    tc = C_BLACK
                    al = PP_ALIGN.LEFT
                elif ci == 1:
                    fill = RGBColor(0xE8, 0xF5, 0xE9)
                    tc = RGBColor(0x1B, 0x5E, 0x20)
                    al = PP_ALIGN.CENTER
                else:
                    fill = RGBColor(0xFF, 0xF8, 0xF8)
                    tc = RGBColor(0x8B, 0x00, 0x00)
                    al = PP_ALIGN.CENTER
                table_cell_set(cell, cell_text, font_size=11, bold=False,
                               text_color=tc, fill_color=fill, align=al)

    add_textbox(sl,
                "First-mover advantage closes in 18–24 months. "
                "The window to lead is now.",
                Cm(1.2), H - Cm(1.9), W - Cm(2.4), Cm(0.8),
                font_size=12, bold=True, italic=True, color=C_BLUE,
                align=PP_ALIGN.CENTER)
    chrome(sl, 18, dark=False)


def slide19(prs):
    """RISK MITIGATION"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "We've Stress-Tested Every Risk")

    risks = [
        ("3D performance on mobile",       "🔴 High",   "CDN + progressive loading + LOD optimisation"),
        ("AI Scoring latency",             "🟡 Medium", "Edge computing + async processing"),
        ("CRM/ERP integration complexity", "🟡 Medium", "API-first + middleware adapter"),
        ("HNW data security",              "🔴 High",   "E2E encryption + PDPA + pen testing"),
        ("Concierge app low adoption",     "🟡 Medium", "White-glove onboarding + incentives"),
        ("Vendor dependency",              "🟡 Medium", "Multi-vendor + IP ownership contract"),
        ("Timeline/scope creep",           "🔴 High",   "Agile sprints + fixed-price phases + CR protocol"),
    ]
    hdr = ["Risk", "Level", "Mitigation"]
    num_rows = len(risks) + 1
    tbl_top  = Cm(2.7)
    tbl_h    = H - tbl_top - Cm(2.2)
    tbl_w    = W - Cm(2.4)
    tbl = sl.shapes.add_table(num_rows, 3,
                               Cm(1.2), tbl_top, tbl_w, tbl_h).table
    col_widths = [Cm(11), Cm(4), Cm(17)]
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ci, h in enumerate(hdr):
        table_cell_set(tbl.cell(0, ci), h, font_size=12, bold=True,
                       text_color=C_WHITE, fill_color=C_BLUE,
                       align=PP_ALIGN.CENTER)

    for ri, (risk, level, mit) in enumerate(risks, start=1):
        fill = RGBColor(0xFF, 0xF0, 0xF0) if "High" in level \
            else RGBColor(0xFF, 0xFD, 0xE7)
        table_cell_set(tbl.cell(ri, 0), risk, font_size=11,
                       text_color=C_BLACK, fill_color=fill,
                       align=PP_ALIGN.LEFT)
        table_cell_set(tbl.cell(ri, 1), level, font_size=11,
                       text_color=C_BLACK, fill_color=fill,
                       align=PP_ALIGN.CENTER)
        table_cell_set(tbl.cell(ri, 2), mit, font_size=11,
                       text_color=C_BLACK, fill_color=fill,
                       align=PP_ALIGN.LEFT)

    chrome(sl, 19, dark=False)


def slide20(prs):
    """PHASED ROADMAP"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "Three Phases. 24 Months. Full Transformation.")

    phases = [
        ("Phase 1  ·  M1 – M8   ·  $3.5M",
         "Foundation",
         "IVA MVP · Admin Portal · CRM Integration",
         C_GOLD, 0, 8),
        ("Phase 2  ·  M9 – M16  ·  $3.5M",
         "Intelligence",
         "AI Scoring · Concierge MVP · Payment Gateway",
         C_BLUE, 8, 8),
        ("Phase 3  ·  M17 – M24 ·  $3.0M",
         "Ecosystem",
         "Full AI · Partner Marketplace · Community",
         RGBColor(0x4A, 0x90, 0xD9), 16, 8),
    ]
    timeline_w = W - Cm(2.4)
    month_w    = timeline_w / 24
    bar_h      = Cm(2.4)
    tl_top     = Cm(3.2)

    # Month ruler
    for m in range(1, 25):
        lft = Cm(1.2) + (m - 1) * month_w
        add_textbox(sl, str(m),
                    lft, tl_top - Cm(0.8), month_w, Cm(0.7),
                    font_size=9, color=C_BLACK, align=PP_ALIGN.CENTER)

    for i, (label, name, detail, col, start, dur) in enumerate(phases):
        lft = Cm(1.2) + start * month_w
        bw  = dur * month_w
        top = tl_top + i * (bar_h + Cm(0.5))
        add_rect(sl, lft, top, bw, bar_h,
                 fill_color=col, line_color=None)
        tc = C_NAVY if col == C_GOLD else C_WHITE
        add_textbox(sl, label,
                    lft + Cm(0.3), top + Cm(0.2), bw - Cm(0.6), Cm(0.9),
                    font_size=12, bold=True, color=tc)
        add_textbox(sl, detail,
                    lft + Cm(0.3), top + Cm(1.2), bw - Cm(0.6), Cm(1.0),
                    font_size=11, italic=True, color=tc)

    # Milestones
    milestones = [
        ("M8\nIVA Launch",         8),
        ("M16\nConcierge Launch",  16),
        ("M24\nFull Platform",     24),
    ]
    ms_top = tl_top + 3 * (bar_h + Cm(0.5)) + Cm(0.4)
    for label, m in milestones:
        lft = Cm(1.2) + (m - 1) * month_w
        add_rect(sl, lft - Cm(0.05), ms_top, Cm(0.1), Cm(1.6),
                 fill_color=C_GOLD, line_color=None)
        add_textbox(sl, label,
                    lft - Cm(1.5), ms_top + Cm(1.7), Cm(3), Cm(1.4),
                    font_size=11, bold=True, color=C_BLUE,
                    align=PP_ALIGN.CENTER)

    chrome(sl, 20, dark=False)


def slide21(prs):
    """KPI NORTH STAR"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "Success Looks Like This at Month 24")

    left_kpis = [
        ("IVA KPIs", True),
        ("Entry Rate",             "5,000+ / month"),
        ("Qualified Lead Rate",    "> 35%"),
        ("Tour Booking Rate",      "> 40%"),
        ("Conversion to Deposit",  "> 6%"),
    ]
    right_kpis = [
        ("Concierge KPIs", True),
        ("Monthly Active Residents", "> 80%"),
        ("AI Satisfaction Score",    "> 4.2 / 5.0"),
        ("Net Promoter Score",       "> 70"),
        ("Partner GMV",              "$2M+ / year"),
    ]

    half_w = W / 2 - Cm(1.5)
    panel_h = Cm(11.0)
    panel_top = Cm(3.2)

    for col_idx, kpis in enumerate([left_kpis, right_kpis]):
        lft = Cm(1.0) + col_idx * (half_w + Cm(1.0))
        add_rect(sl, lft, panel_top, half_w, panel_h,
                 fill_color=C_BLUE,
                 line_color=C_GOLD, line_width=Pt(1.5))
        for ki, item in enumerate(kpis):
            top = panel_top + Cm(0.4) + ki * Cm(2.0)
            if item[1] is True:
                add_textbox(sl, item[0],
                            lft + Cm(0.5), top, half_w - Cm(1.0), Cm(1.3),
                            font_size=16, bold=True, color=C_GOLD)
            else:
                add_textbox(sl, item[0],
                            lft + Cm(0.5), top, half_w * 0.58, Cm(1.4),
                            font_size=14, color=C_LBLUE)
                add_textbox(sl, item[1],
                            lft + half_w * 0.6, top, half_w * 0.36, Cm(1.4),
                            font_size=14, bold=True, color=C_WHITE,
                            align=PP_ALIGN.RIGHT)

    chrome(sl, 21, dark=True)


def slide22(prs):
    """GOVERNANCE MODEL"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_LIGHT)
    section_title(sl, "Your Platform. Our Expertise. One Team.")

    # Org hierarchy
    nodes = [
        ("Steering Committee\n(Client C-Suite)",   Cm(11),  Cm(2.8),  Cm(11), Cm(1.8), C_GOLD,  C_NAVY),
        ("→ Monthly Review",                        Cm(22.5),Cm(3.0),  Cm(7.5),Cm(1.4), None,     C_BLUE),
        ("Product Owner\n(Client Representative)", Cm(11),  Cm(5.4),  Cm(11), Cm(1.8), C_BLUE,   C_WHITE),
        ("→ Sprint Review & UAT",                   Cm(22.5),Cm(5.6),  Cm(7.5),Cm(1.4), None,     C_BLUE),
        ("Delivery Lead  +  Solution Architect\n(Vendor)",
         Cm(8.5), Cm(8.0), Cm(16), Cm(1.8), RGBColor(0x14, 0x2A, 0x42), C_WHITE),
    ]
    for (text, lft, top, wid, hei, fill, tc) in nodes:
        if fill:
            add_rect(sl, lft, top, wid, hei,
                     fill_color=fill, line_color=C_GOLD, line_width=Pt(1))
            add_textbox(sl, text, lft + Cm(0.3), top + Cm(0.1),
                        wid - Cm(0.6), hei - Cm(0.2),
                        font_size=13, bold=True, color=tc,
                        align=PP_ALIGN.CENTER)
        else:
            add_textbox(sl, text, lft, top, wid, hei,
                        font_size=12, italic=True, color=tc)

    # Connector lines
    for y_from, y_to in [(Cm(4.6), Cm(5.4)), (Cm(7.2), Cm(8.0))]:
        ln = sl.shapes.add_connector(1, W/2, y_from, W/2, y_to)
        ln.line.color.rgb = C_GOLD
        ln.line.width = Pt(1.5)

    squads = [
        "IVA / 3D / VR", "AI / Data", "Concierge App",
        "Integrations", "Admin / Ops", "QA / Security",
    ]
    sq_w = (W - Cm(2.4)) / 6
    sq_top = Cm(10.5)
    for i, sq in enumerate(squads):
        lft = Cm(1.2) + i * sq_w
        add_rect(sl, lft + Cm(0.1), sq_top, sq_w - Cm(0.2), Cm(2.2),
                 fill_color=C_BLUE,
                 line_color=C_GOLD, line_width=Pt(1))
        add_textbox(sl, sq, lft + Cm(0.2), sq_top + Cm(0.4),
                    sq_w - Cm(0.4), Cm(1.4),
                    font_size=11, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)

    add_textbox(sl,
                "2-week Sprints  ·  Monthly Milestone Review  ·  "
                "Quarterly Strategic Review",
                Cm(1.2), H - Cm(2.5), W - Cm(2.4), Cm(0.9),
                font_size=12, bold=True, color=C_BLUE,
                align=PP_ALIGN.CENTER)
    chrome(sl, 22, dark=False)


def slide23(prs):
    """WHY US"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "We've Done This Before. At Scale.")

    refs = [
        ("[Reference Project 1]",
         "Luxury Real Estate Platform",
         "Full pre-sale digital platform · 3D tours · AI lead scoring\n"
         "Delivered in 14 months · 42% conversion uplift"),
        ("[Reference Project 2]",
         "AI Concierge Implementation",
         "Residential app · 500+ services · 4.6/5 satisfaction\n"
         "12,000 active residents · 78% daily engagement"),
        ("[Reference Project 3]",
         "Metaverse Showroom",
         "Web-based 3D showroom · 6 unit types · Avatar tours\n"
         "120,000 unique visitors in launch month"),
    ]
    card_w = Cm(9.8)
    card_h = Cm(6.8)
    top    = Cm(3.0)
    for i, (code, title, body) in enumerate(refs):
        lft = Cm(1.2) + i * (card_w + Cm(0.55))
        add_rect(sl, lft, top, card_w, card_h,
                 fill_color=C_BLUE,
                 line_color=C_GOLD, line_width=Pt(2))
        add_textbox(sl, code,
                    lft + Cm(0.4), top + Cm(0.3), card_w - Cm(0.8), Cm(1.1),
                    font_size=13, italic=True, color=C_GOLD)
        add_textbox(sl, title,
                    lft + Cm(0.4), top + Cm(1.5), card_w - Cm(0.8), Cm(1.2),
                    font_size=14, bold=True, color=C_WHITE)
        add_textbox(sl, body,
                    lft + Cm(0.4), top + Cm(3.0), card_w - Cm(0.8), Cm(3.4),
                    font_size=12, color=C_LBLUE)

    add_textbox(sl,
                "Tech Partners:  AWS  ·  Salesforce  ·  Three.js  ·  OpenAI  ·  Flutter",
                Cm(1.2), Cm(10.8), W - Cm(2.4), Cm(0.9),
                font_size=13, bold=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)
    creds = [
        "▸  Lead architect: 12+ years in PropTech / LuxTech platforms",
        "▸  AI team: published research in behavioural scoring & NLP",
        "▸  Delivery track record: 23 enterprise projects, 0 missed go-live dates",
    ]
    for i, cred in enumerate(creds):
        add_textbox(sl, cred,
                    Cm(2.0), Cm(12.0) + i * Cm(1.3), W - Cm(4), Cm(1.1),
                    font_size=13, color=C_LBLUE)

    chrome(sl, 23, dark=True)


def slide24(prs):
    """NEXT STEPS"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)
    section_title(sl, "From Proposal to Platform in 6 Steps")

    steps = [
        ("1", "Alignment Workshop",       "Week 1",    "Confirm scope & priorities"),
        ("2", "Technical Discovery",      "Week 2–3",  "Architecture validation"),
        ("3", "Contract & Vendor Select", "Week 4–5",  "RFP / direct award"),
        ("4", "Kickoff & Phase 1 Start",  "Week 6",    "Teams formed, backlog ready"),
        ("5", "Phase 1 Demo",             "Month 4",   "First IVA preview to stakeholders"),
        ("6", "Phase 1 Go-Live",          "Month 8",   "Pre-sale campaign launch"),
    ]
    step_w  = Cm(5.2)
    step_h  = Cm(5.8)
    row_top = Cm(3.0)
    for i, (num, title, timing, desc) in enumerate(steps):
        lft = Cm(0.8) + i * (step_w + Cm(0.4))
        # Gold circle (rect)
        add_rect(sl, lft + step_w / 2 - Cm(0.8), row_top,
                 Cm(1.6), Cm(1.6),
                 fill_color=C_GOLD, line_color=None)
        add_textbox(sl, num,
                    lft + step_w / 2 - Cm(0.8),
                    row_top + Cm(0.1),
                    Cm(1.6), Cm(1.4),
                    font_size=18, bold=True, color=C_NAVY,
                    align=PP_ALIGN.CENTER)
        add_rect(sl, lft, row_top + Cm(1.8), step_w, step_h - Cm(1.8),
                 fill_color=C_BLUE,
                 line_color=C_GOLD, line_width=Pt(1))
        add_textbox(sl, title,
                    lft + Cm(0.2), row_top + Cm(2.0),
                    step_w - Cm(0.4), Cm(1.4),
                    font_size=13, bold=True, color=C_GOLD,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, timing,
                    lft + Cm(0.2), row_top + Cm(3.5),
                    step_w - Cm(0.4), Cm(0.9),
                    font_size=12, italic=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, desc,
                    lft + Cm(0.2), row_top + Cm(4.5),
                    step_w - Cm(0.4), Cm(1.0),
                    font_size=11, color=C_LBLUE,
                    align=PP_ALIGN.CENTER)

    # CTA
    cta_top = H - Cm(3.0)
    add_rect(sl, Cm(5.5), cta_top, W - Cm(11), Cm(1.4),
             fill_color=C_GOLD, line_color=None)
    add_textbox(sl,
                "Decision needed by [Date] to hit pre-sale launch window",
                Cm(5.7), cta_top + Cm(0.1), W - Cm(11.4), Cm(1.2),
                font_size=15, bold=True, color=C_NAVY,
                align=PP_ALIGN.CENTER)

    chrome(sl, 24, dark=True)


def slide25(prs):
    """CLOSING"""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl, C_NAVY)

    add_textbox(sl, "This Is More Than a Platform.",
                Cm(2), Cm(1.8), W - Cm(4), Cm(2.0),
                font_size=34, bold=True, italic=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)

    body = (
        "It is the first impression every future resident will have of Splendora HH4 —\n"
        "before they visit, before they buy, before they move in.\n\n"
        "It is the daily companion they will trust for years after.\n\n"
        "It is the data, the relationship, and the brand equity\n"
        "that no competitor can replicate overnight.\n\n"
        "The question is not whether to build it.\n"
        "The question is whether to build it first."
    )
    add_textbox(sl, body,
                Cm(3), Cm(4.2), W - Cm(6), Cm(10.0),
                font_size=18, italic=True, color=C_LBLUE,
                align=PP_ALIGN.CENTER)

    # Gold line
    line = sl.shapes.add_connector(1,
                                   Cm(4), H - Cm(2.8),
                                   W - Cm(4), H - Cm(2.8))
    line.line.color.rgb = C_GOLD
    line.line.width = Pt(2)

    add_textbox(sl,
                "HH4 Digital Platform  ·  Confidential  ·  March 2026",
                Cm(2), H - Cm(2.3), W - Cm(4), Cm(0.9),
                font_size=13, italic=True, color=C_GOLD,
                align=PP_ALIGN.CENTER)

    add_logo(sl, dark=True)
    add_slide_number(sl, 25, dark=True)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    builders = [
        slide01, slide02, slide03, slide04, slide05,
        slide06, slide07, slide08, slide09, slide10,
        slide11, slide12, slide13, slide14, slide15,
        slide16, slide17, slide18, slide19, slide20,
        slide21, slide22, slide23, slide24, slide25,
    ]
    for fn in builders:
        fn(prs)

    out_dir  = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "HH4_Platform_Presentation.pptx")
    prs.save(out_path)
    print(f"✅  Saved: {out_path}")


if __name__ == "__main__":
    main()
